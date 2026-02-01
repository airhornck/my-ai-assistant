"""
文档解析：从文件提取可读文本。
支持 PDF、TXT、DOCX、PPTX、MD、图片（OCR）。生产环境可扩展更多格式。
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# 支持的文档格式
SUPPORTED_DOC_EXTENSIONS = frozenset({
    "pdf", "txt", "text", "md", "markdown", "docx", "doc",
    "pptx", "ppt",
    "jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "tif",
})


def _guess_file_type(filename: str) -> str:
    ext = (filename or "").rsplit(".", 1)[-1].lower() if "." in (filename or "") else ""
    return ext or "bin"


def parse_text(storage_path: str, file_type: str, filename: str = "") -> str:
    """
    从存储路径解析文档内容为纯文本。
    支持 pdf, txt, docx, pptx, md, 图片(jpg/png等)。
    路径支持相对路径（相对当前工作目录）或绝对路径。
    """
    import os
    path = Path(storage_path)
    if not path.is_absolute():
        path = Path(os.path.abspath(storage_path))
    if not path.exists() or not path.is_file():
        logger.warning("文档不存在或非文件: %s (resolved: %s)", storage_path, path)
        return ""
    ft = file_type.lower() or _guess_file_type(filename or path.name)
    resolved = str(path)
    try:
        if ft == "pdf":
            return _parse_pdf(resolved)
        if ft in ("txt", "text", "md", "markdown"):
            return _parse_txt(resolved)
        if ft in ("docx", "doc"):
            return _parse_docx(resolved)
        if ft in ("pptx", "ppt"):
            return _parse_pptx(resolved)
        if ft in ("jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "tif"):
            return _parse_image(resolved)
    except Exception as e:
        logger.warning("文档解析失败 %s: %s", storage_path, e, exc_info=True)
    return ""


def _parse_txt(path: str) -> str:
    """解析 TXT/MD 等纯文本。"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read().strip()


def _parse_pdf(path: str) -> str:
    """解析 PDF。"""
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf 未安装，无法解析 PDF")
        return ""
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        try:
            t = page.extract_text()
            if t:
                parts.append(t.strip())
        except Exception:
            pass
    return "\n\n".join(parts)


def _parse_docx(path: str) -> str:
    """解析 DOCX。"""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        logger.warning("python-docx 未安装，无法解析 DOCX")
        return ""
    doc = DocxDocument(path)
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def _parse_pptx(path: str) -> str:
    """解析 PPTX（PowerPoint）。仅支持 .pptx，.ppt 旧格式需先转换。"""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx 未安装，无法解析 PPTX")
        return ""
    ext = Path(path).suffix.lower()
    if ext != ".pptx":
        logger.warning("仅支持 .pptx 格式，.ppt 需先转换为 .pptx")
        return ""
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text.strip()
                if t:
                    parts.append(t)
    return "\n\n".join(parts)


def _parse_image(path: str) -> str:
    """解析图片：OCR 提取文字。需安装 pytesseract 和 Pillow，系统需安装 Tesseract-OCR。"""
    try:
        import pytesseract
        from PIL import Image
    except ImportError as e:
        logger.warning("pytesseract 或 Pillow 未安装，无法解析图片: %s", e)
        return ""
    try:
        img = Image.open(path)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        return (text or "").strip()
    except Exception as e:
        logger.warning("图片 OCR 失败 %s: %s", path, e)
        return ""


class DocumentParser:
    """
    文档解析器：对外统一接口。
    """

    def parse(self, storage_path: str, file_type: str, filename: str = "") -> str:
        return parse_text(storage_path, file_type, filename)
